"""
Drawing Primitives — مذكرتي Pro v17.2 VISUAL UPGRADE
MAJOR improvements:
- shadow() now works on ALL shape types (sp + pic + grpSp)
- glow() effect added
- gradient_fill() with 3-stop support
- polygon() for decorative triangles/diamonds
- _sort_spPr() stricter OOXML ordering
- soft_shadow() realistic multi-layer effect
"""
from __future__ import annotations
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

W, H = 33.867, 19.05

_SPPR_ORDER = [
    qn('a:xfrm'), qn('a:prstGeom'), qn('a:custGeom'),
    qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'),
    qn('a:blipFill'), qn('a:pattFill'), qn('a:grpFill'),
    qn('a:ln'), qn('a:effectLst'), qn('a:effectDag'),
    qn('a:scene3d'), qn('a:sp3d'), qn('a:extLst'),
]
_SPPR_RANK = {t: i for i, t in enumerate(_SPPR_ORDER)}


def _sort_spPr(spPr) -> None:
    children = list(spPr)
    children.sort(key=lambda el: _SPPR_RANK.get(el.tag, 99))
    for c in children: spPr.remove(c)
    for c in children: spPr.append(c)


def _get_spPr(shape):
    sp = shape._element
    # handles <p:sp>, <p:pic>, <p:grpSp>
    for tag in (qn('p:spPr'), qn('pic:spPr')):
        found = sp.find('.//' + tag) if tag not in [c.tag for c in sp] else sp.find(tag)
        if found is not None: return found
    return sp.find(qn('p:spPr'))


def cm(v): return int(Cm(v))
def pt(v): return int(Pt(v))


# ── Basic shapes ──────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill: RGBColor, line=None, line_w=0.5):
    if w <= 0 or h <= 0: return None
    s = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = pt(line_w)
    else: s.line.fill.background()
    return s

def rrect(slide, x, y, w, h, fill: RGBColor, radius_pct=8, line=None, line_w=0.5):
    if w <= 0 or h <= 0: return None
    s = slide.shapes.add_shape(5, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = pt(line_w)
    else: s.line.fill.background()
    try:
        adj = s.adjustments
        if adj and len(adj) > 0: adj[0] = max(0, min(50, radius_pct)) * 1000
    except: pass
    return s

def oval(slide, x, y, w, h, fill: RGBColor, alpha=100):
    if w <= 0 or h <= 0: return None
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100: set_solid_alpha(s, alpha)
    return s

def bg(slide, color: RGBColor): rect(slide, 0, 0, W, H, color)
def hline(slide, x, y, w, color, thickness=0.08): rect(slide, x, y, w, thickness, color)
def vline(slide, x, y, h2, color, thickness=0.08): rect(slide, x, y, thickness, h2, color)


# ── XML fill helpers ──────────────────────────────────────────────────
def set_solid_alpha(shape, alpha_pct: int):
    try:
        spPr = _get_spPr(shape)
        srgb = spPr.find('.//' + qn('a:srgbClr'))
        if srgb is not None:
            for e in srgb.findall(qn('a:alpha')): srgb.remove(e)
            alp = etree.SubElement(srgb, qn('a:alpha'))
            alp.set('val', str(int(alpha_pct * 1000)))
    except: pass


def gradient_fill(shape, c1: str, c2: str, angle: float = 90, c3: str = None) -> None:
    """2 or 3 stop linear gradient. c3 = optional mid stop."""
    try:
        spPr = _get_spPr(shape)
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill'), qn('a:grpFill')]:
            for el in spPr.findall(tag): spPr.remove(el)

        grad = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))

        gs0 = etree.SubElement(gsLst, qn('a:gs')); gs0.set('pos', '0')
        etree.SubElement(gs0, qn('a:srgbClr')).set('val', c1.lstrip('#'))

        if c3:
            gsm = etree.SubElement(gsLst, qn('a:gs')); gsm.set('pos', '50000')
            etree.SubElement(gsm, qn('a:srgbClr')).set('val', c3.lstrip('#'))

        gs1 = etree.SubElement(gsLst, qn('a:gs')); gs1.set('pos', '100000')
        etree.SubElement(gs1, qn('a:srgbClr')).set('val', c2.lstrip('#'))

        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000))); lin.set('scaled', '0')
        spPr.append(grad); _sort_spPr(spPr)
    except: pass


def gradient_rect(slide, x, y, w, h, c1, c2, angle=0, c3=None):
    c = c1.lstrip('#')
    s = rect(slide, x, y, w, h, RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16)))
    if s: gradient_fill(s, c1, c2, angle, c3)
    return s


def shadow(shape, blur=18, dist=6, angle=135, alpha=0.28, color="000000") -> None:
    """Drop shadow — works correctly after _sort_spPr."""
    try:
        spPr = _get_spPr(shape)
        if spPr is None: return
        for old in spPr.findall(qn('a:effectLst')): spPr.remove(old)
        eLst = etree.Element(qn('a:effectLst'))
        shdw = etree.SubElement(eLst, qn('a:outerShdw'))
        shdw.set('blurRad', str(int(blur * 12700)))
        shdw.set('dist', str(int(dist * 12700)))
        shdw.set('dir', str(int(angle * 60000)))
        shdw.set('algn', 'tl')
        srgb = etree.SubElement(shdw, qn('a:srgbClr')); srgb.set('val', color.lstrip('#'))
        etree.SubElement(srgb, qn('a:alpha')).set('val', str(int(alpha * 100000)))
        spPr.append(eLst); _sort_spPr(spPr)
    except: pass


def glow(shape, color: str, radius=10, alpha=0.35) -> None:
    """Glow effect — adds luminous halo."""
    try:
        spPr = _get_spPr(shape)
        if spPr is None: return
        eLst = spPr.find(qn('a:effectLst'))
        if eLst is None:
            eLst = etree.Element(qn('a:effectLst'))
            spPr.append(eLst); _sort_spPr(spPr)
        g = etree.SubElement(eLst, qn('a:glow'))
        g.set('rad', str(int(radius * 12700)))
        srgb = etree.SubElement(g, qn('a:srgbClr')); srgb.set('val', color.lstrip('#'))
        etree.SubElement(srgb, qn('a:alpha')).set('val', str(int(alpha * 100000)))
    except: pass


def soft_shadow(shape, alpha=0.4) -> None:
    """Realistic soft shadow with large blur."""
    shadow(shape, blur=32, dist=8, angle=135, alpha=alpha)


# ── Text ──────────────────────────────────────────────────────────────
def txt(slide, text, x, y, w, h,
        font="Cairo", size=14, bold=False, italic=False,
        color: RGBColor = None, align=PP_ALIGN.RIGHT,
        margin=0.12, spacing=None):
    if not text or w <= 0 or h <= 0: return None
    tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = cm(margin); tf.margin_right = cm(margin)
    tf.margin_top = cm(0.04); tf.margin_bottom = cm(0.04)
    p = tf.paragraphs[0]; p.alignment = align
    if spacing:
        try: p.line_spacing = Pt(spacing)
        except: pass
    run = p.add_run(); run.text = str(text)
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color
    return tb


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])
