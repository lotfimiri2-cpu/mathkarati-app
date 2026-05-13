"""
مذكرتي Pro — ULTRA CANVA ENGINE v2
====================================
تصاميم سينمائية بمستوى Behance / Dribbble
- خطوط متدرجة الحجم والوزن بهرمية واضحة
- أشكال هندسية معقدة (مثلثات، شرائح، خطوط مائلة)
- تدرجات لونية متطورة
- نظام بطاقات متعدد المستويات مع ظلال حقيقية
- 3 عائلات بصرية مختلفة جذرياً: NOIR · VIVID · MINIMAL
"""

import sys, json, math, datetime
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

W, H   = 13.33, 7.50
MX, MY = 0.60, 0.48

# ── خطوط بديلة مضمونة عند غياب الخط الأصلي ───────────────────────
# الخطوط المضمونة على أي نظام (Windows/Linux/Mac/Render)
_GUARANTEED = {"Arial", "Calibri", "Tahoma", "Times New Roman", "Courier New"}
_FONT_FALLBACK = {
    "Palatino Linotype": "Georgia",
    "Trebuchet MS":      "Arial",
    "Segoe UI Emoji":    "Arial",
    "Segoe UI":          "Arial",
}

def _cairo_available() -> bool:
    """يتحقق من وجود خط Cairo فعلياً على النظام عبر fontconfig أو مسارات مباشرة."""
    import shutil, os
    # 1) fc-list (Linux/Render)
    if shutil.which("fc-list"):
        try:
            import subprocess
            out = subprocess.run(["fc-list", ":family=Cairo"],
                                 capture_output=True, text=True, timeout=3)
            if "cairo" in out.stdout.lower():
                return True
        except Exception:
            pass
    # 2) مسارات شائعة لملفات الخط
    search_dirs = [
        "/usr/share/fonts", "/usr/local/share/fonts",
        os.path.expanduser("~/.fonts"),
        "C:/Windows/Fonts",
        "/Library/Fonts", os.path.expanduser("~/Library/Fonts"),
    ]
    for d in search_dirs:
        if os.path.isdir(d):
            for root, _, files in os.walk(d):
                for f in files:
                    if "cairo" in f.lower() and f.lower().endswith((".ttf", ".otf")):
                        return True
    return False

# التحقق مرة واحدة عند الاستيراد
_CAIRO_OK = _cairo_available()
_FONT_FALLBACK["Cairo"] = "Cairo" if _CAIRO_OK else "Arial"

if not _CAIRO_OK:
    import logging as _log
    _log.getLogger(__name__).warning(
        "⚠️  خط Cairo غير موجود على النظام — سيُستخدم Arial بديلاً. "
        "يُنصح بتشغيل build.sh لتثبيته."
    )

def safe_font(name: str) -> str:
    """يُعيد الخط إذا كان آمناً أو بديله المضمون."""
    if name in _GUARANTEED:
        return name
    return _FONT_FALLBACK.get(name, "Arial")

def rgb(r,g,b): return RGBColor(r,g,b)
def hx(h):      return RGBColor.from_string(h.lstrip('#'))
def safe(v,fb=""): return str(v).strip() if v else fb
def clamp(v,lo,hi): return max(lo,min(hi,v))
def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])
def cm(v): return Cm(v)
def emu(v): return int(Cm(v))

# ─── Core Drawing ─────────────────────────────────────────────────
def rect(slide, x,y,w,h, fill, line_color=None, line_w=0.6):
    if w<=0 or h<=0: return None
    s = slide.shapes.add_shape(1, cm(x),cm(y),cm(w),cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color: s.line.color.rgb=line_color; s.line.width=Pt(line_w)
    else: s.line.fill.background()
    return s

def rrect(slide, x,y,w,h, fill, r_pct=10, line_color=None, line_w=0.5):
    if w<=0 or h<=0: return None
    s = slide.shapes.add_shape(5, cm(x),cm(y),cm(w),cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color: s.line.color.rgb=line_color; s.line.width=Pt(line_w)
    else: s.line.fill.background()
    try:
        adj = s.adjustments
        if adj and len(adj)>0: adj[0] = clamp(r_pct,0,50)*1000
    except: pass
    return s

def oval(slide, x,y,w,h, fill):
    if w<=0 or h<=0: return None
    s = slide.shapes.add_shape(9, cm(x),cm(y),cm(w),cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def triangle(slide, x,y,w,h, fill, flip_h=False):
    """مثلث قائم الزاوية"""
    if w<=0 or h<=0: return None
    shape_id = 6 if not flip_h else 7  # rightTriangle
    try:
        s = slide.shapes.add_shape(6, cm(x),cm(y),cm(w),cm(h))
        s.fill.solid(); s.fill.fore_color.rgb = fill
        s.line.fill.background()
        return s
    except: return None

def parallelogram(slide, x,y,w,h, fill, slant=15):
    """متوازي أضلاع — يُحاكى بمستطيل مائل"""
    try:
        s = slide.shapes.add_shape(8, cm(x),cm(y),cm(w),cm(h))
        s.fill.solid(); s.fill.fore_color.rgb = fill
        s.line.fill.background()
        return s
    except: return None

def bg(slide, color): rect(slide,0,0,W,H,color)

def lh(slide, x,y,w, color, h=0.045): rect(slide,x,y,w,h,color)
def lv(slide, x,y,h2, color, w=0.045): rect(slide,x,y,w,h2,color)

def shadow_xml(shape, blur=10, dist=4, angle=135, alpha=0.18):
    """ظل حقيقي عبر XML"""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None: return
        # Remove existing effectLst
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        eLst = etree.SubElement(spPr, qn('a:effectLst'))
        shdw = etree.SubElement(eLst, qn('a:outerShdw'))
        shdw.set('blurRad', str(int(blur*12700)))
        shdw.set('dist',    str(int(dist*12700)))
        shdw.set('dir',     str(int(angle*60000)))
        shdw.set('algn',    'tl')
        srgb = etree.SubElement(shdw, qn('a:srgbClr'))
        srgb.set('val', '000000')
        alp  = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha*100000)))
    except: pass

def txt(slide, text, x,y,w,h,
        font="Cairo", size=13, bold=False, italic=False,
        color=None, align=PP_ALIGN.RIGHT, mg=0.07,
        rtl=True, valign="top", spacing=None):
    if w<=0 or h<=0 or not text: return None
    tb = slide.shapes.add_textbox(cm(x),cm(y),cm(w),cm(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=cm(mg); tf.margin_right=cm(mg)
    tf.margin_top=cm(0.03); tf.margin_bottom=cm(0.03)
    p = tf.paragraphs[0]; p.alignment = align
    if spacing:
        try:
            from pptx.util import Pt as _Pt
            p.line_spacing = _Pt(spacing)
        except: pass
    run = p.add_run()
    run.text = str(text)
    run.font.name = safe_font(font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    return tb

# ═══════════════════════════════════════════════════════════════════
# TYPOGRAPHY SYSTEM — هرمية الخطوط الكاملة
# ═══════════════════════════════════════════════════════════════════
class TY:
    """
    هرمية الخطوط:
    DISPLAY  52pt Bold    — عنوان الغلاف الرئيسي
    H1       30pt Bold    — عنوان الشريحة
    H2       18pt Bold    — عنوان قسم / بطاقة
    H3       14pt Bold    — عنوان فرعي داخلي
    H4       12pt Bold    — تسمية عنصر
    BODY     13pt Regular — فقرة نصية رئيسية
    BODY_SM  11pt Regular — فقرة ثانوية
    BULLET   12pt Regular — عنصر قائمة
    CAPTION   9pt Italic  — تعليق / مرجع
    LABEL    13pt Bold    — تسمية مربع / خلية
    META     11pt Regular — معلومات هوية
    NUMBER   32pt Bold    — رقم ترتيبي
    KPI_LG   52pt Bold    — قيمة KPI كبيرة
    OVER      9pt Bold    — overline (أحرف كبيرة)
    """
    DISPLAY = (52, True,  False, "Georgia")
    H1      = (30, True,  False, None)     # HF from palette
    H2      = (18, True,  False, None)
    H3      = (14, True,  False, None)
    H4      = (12, True,  False, None)
    BODY    = (13, False, False, None)
    BODY_SM = (11, False, False, None)
    BULLET  = (12, False, False, None)
    CAPTION = ( 9, False, True,  "Calibri")
    LABEL   = (13, True,  False, None)
    META    = (11, False, False, None)
    NUMBER  = (32, True,  False, "Calibri")
    KPI_LG  = (52, True,  False, "Calibri")
    OVER    = ( 9, True,  False, "Calibri")

def t(slide, text, x,y,w,h, ty, T, color=None, align=PP_ALIGN.RIGHT, light=True, custom_font=None):
    """Typography helper — uses TY constants + palette"""
    size, bold, italic, force_font = ty
    font = custom_font or force_font or (T["HF"] if size>=18 else T["BF"])
    if color is None:
        color = T["TL"] if light else T["TD"]
    return txt(slide, text, x,y,w,h,
               font=font, size=size, bold=bold, italic=italic,
               color=color, align=align)

def overline(slide, text, x,y,w, T, color=None, align=PP_ALIGN.LEFT):
    c = color or T["A"]
    tb = txt(slide, text.upper(), x,y,w,0.30,
             font="Calibri", size=TY.OVER[0], bold=True,
             color=c, align=align, rtl=False)
    return tb

# ═══════════════════════════════════════════════════════════════════
# COLOR SYSTEM — 8 لوحات + 3 عائلات
# ═══════════════════════════════════════════════════════════════════
def PAL(d,m,l, a,a2, tl,td,tm, g1,g2, sc_list, hf,bf, fam,
        card_d, card_l, border):
    return {
        "D":hx(d), "M":hx(m), "L":hx(l),
        "A":hx(a), "A2":hx(a2),
        "TL":hx(tl), "TD":hx(td), "TM":hx(tm),
        "G1":hx(g1), "G2":hx(g2),
        "SC":[hx(c) for c in sc_list],
        "HF":safe_font(hf), "BF":safe_font(bf), "FAM":fam,
        "CD":hx(card_d), "CL":hx(card_l),
        "BO":hx(border),
        "CB":hx("FFFFFF"), "CE":hx("E8EDF2"),
    }

PALETTES = {
# ── NOIR family — فاخر داكن ──────────────────────────────────────
"navy_gold": PAL(
    "050F1E","0B1D38","EEF3FB",
    "D4A017","F5D060",
    "FFFFFF","050F1E","7A90AA",
    "0B1D38","D4A017",
    ["D4A017","4DBFFF","F5D060","FF6B6B","00D4AA","BF7FFF","FF9D3A"],
    "Palatino Linotype","Cairo","NOIR",
    "0B1D38","FFFFFF","1A3A5C"),

"midnight_purple": PAL(
    "0A0418","1E0852","F4F0FF",
    "C060FF","E0A8FF",
    "FFFFFF","0A0418","8A7AAA",
    "1E0852","C060FF",
    ["C060FF","FF5C8A","E0A8FF","FFD060","00D4AA","4DBFFF","FF7A40"],
    "Georgia","Cairo","NOIR",
    "1E0852","FFFFFF","3A1A88"),

"forest": PAL(
    "051810","0F3020","EDFAF0",
    "5ABF40","A8E890",
    "FFFFFF","051810","508060",
    "0F3020","5ABF40",
    ["5ABF40","4DBFFF","A8E890","FFD060","FF6B6B","C060FF","FF7A40"],
    "Georgia","Cairo","NOIR",
    "0F3020","FFFFFF","1A5030"),

"sand_gold": PAL(
    "1E0E04","4A2808","FFF8EE",
    "D08820","F0C060",
    "FFFFFF","1E0E04","8A7040",
    "4A2808","D08820",
    ["D08820","4DBFFF","F0C060","FF6B6B","00D4AA","C060FF","5ABF40"],
    "Palatino Linotype","Cairo","NOIR",
    "4A2808","FFFFFF","704010"),

# ── VIVID family — حيوي عصري ─────────────────────────────────────
"dark_teal": PAL(
    "071828","0E3850","E4F6FF",
    "00BFA0","60EDD8",
    "FFFFFF","071828","3A7888",
    "0E3850","00BFA0",
    ["00BFA0","FF5C35","60EDD8","FFD060","C060FF","4DBFFF","FF6B6B"],
    "Trebuchet MS","Cairo","VIVID",
    "0E3850","FFFFFF","1A5870"),

"charcoal_orange": PAL(
    "141420","28284A","FFF5F0",
    "FF5820","FF9060",
    "FFFFFF","141420","7070A0",
    "28284A","FF5820",
    ["FF5820","00BFA0","FF9060","FFD060","C060FF","4DBFFF","5ABF40"],
    "Trebuchet MS","Cairo","VIVID",
    "28284A","FFFFFF","404070"),

"burgundy": PAL(
    "200A18","501030","FFF0F6",
    "E82880","FFB0D0",
    "FFFFFF","200A18","904060",
    "501030","E82880",
    ["E82880","4DBFFF","FFB0D0","FFD060","00BFA0","C060FF","FF5820"],
    "Georgia","Cairo","VIVID",
    "501030","FFFFFF","801848"),

# ── MINIMAL family — نظيف احترافي ────────────────────────────────
"ice_blue": PAL(
    "061840","1040A0","EAF2FF",
    "1060E0","70B0FF",
    "FFFFFF","061840","3060A0",
    "061840","1060E0",
    ["1060E0","FF5820","70B0FF","FFD060","00BFA0","C060FF","5ABF40"],
    "Calibri","Cairo","MINIMAL",
    "061840","FFFFFF","1840A0"),

}


# ═══════════════════════════════════════════════════════════════════
# ADVANCED DECORATORS
# ═══════════════════════════════════════════════════════════════════

def deco_circles(slide, T, configs):
    """دوائر ديكورية متعددة — كل واحدة: (cx,cy,r,color,alpha_pct)"""
    for cx,cy,r,color,alpha in configs:
        s = oval(slide, cx-r, cy-r, r*2, r*2, color)
        if s and alpha < 100:
            try:
                # Apply transparency via XML
                sp = s._element
                spPr = sp.find(qn('p:spPr'))
                fld = spPr.find('.//' + qn('a:solidFill'))
                if fld is not None:
                    srgb = fld.find(qn('a:srgbClr'))
                    if srgb is not None:
                        alpha_e = etree.SubElement(srgb, qn('a:alpha'))
                        alpha_e.set('val', str(alpha*1000))
            except: pass

def deco_diagonal_band(slide, T, x,y,w,h, color, angle_deg=15, n=4, spacing=0.8):
    """شرائح مائلة ديكورية"""
    for i in range(n):
        xi = x + i*spacing
        rect(slide, xi, y, 0.08, h, color)

def deco_corner_accent(slide, T, corner='br', size=2.5):
    """زاوية مزخرفة بأشكال هندسية"""
    c = T["A"]
    if corner == 'br':
        oval(slide, W-size*1.2, H-size*1.2, size*1.8, size*1.8, c)
        oval(slide, W-size*0.7, H-size*0.7, size*1.0, size*1.0, T["M"])
    elif corner == 'tl':
        oval(slide, -size*0.6, -size*0.6, size*1.8, size*1.8, c)

def gradient_bar(slide, x,y,w,h, colors, vertical=False):
    """شريط تدرج لوني — يُحاكى بمستطيلات متعددة"""
    n = len(colors)
    for i,c in enumerate(colors):
        if vertical:
            rect(slide, x, y+i*(h/n), w, h/n+0.02, c)
        else:
            rect(slide, x+i*(w/n), y, w/n+0.02, h, c)

def decorative_grid(slide, x,y,w,h, color, rows=3, cols=4):
    """شبكة نقاط ديكورية"""
    gx = w/cols; gy = h/rows
    for r in range(rows):
        for c in range(cols):
            dx = x + c*gx + gx/2 - 0.04
            dy = y + r*gy + gy/2 - 0.04
            oval(slide, dx, dy, 0.08, 0.08, color)

def accent_shape(slide, T, shape='arc', x=0, y=0, w=2, h=2):
    """شكل accent متقدم"""
    # نستخدم مجموعة أشكال لتوليد تأثيرات بصرية
    oval(slide, x, y, w, h, T["A"])
    oval(slide, x+w*0.15, y+h*0.15, w*0.70, h*0.70, T["D"])

def card_premium(slide, x,y,w,h, T, accent, style='dark'):
    """بطاقة بمستوى premium مع كل التأثيرات"""
    bg_c = T["CD"] if style=='dark' else T["CB"]
    s = rrect(slide, x,y,w,h, bg_c, r_pct=8,
              line_color=T["BO"] if style=='light' else None)
    if s: shadow_xml(s, blur=12, dist=4, alpha=0.16)
    # شريط علوي
    rect(slide, x, y, w, 0.09, accent)
    # شريط جانبي يمين
    rect(slide, x, y, 0.10, h, accent)
    return s

def kpi_premium(slide, x,y,w,h, T, value, label, color, sub=None):
    """KPI بتصميم Ultra Premium"""
    # ظل
    s_bg = rrect(slide, x+0.06, y+0.06, w, h, hx("000000"), r_pct=8)
    if s_bg:
        try:
            sp = s_bg._element
            spPr = sp.find(qn('p:spPr'))
            fld = spPr.find('.//' + qn('a:solidFill'))
            if fld is not None:
                srgb = fld.find(qn('a:srgbClr'))
                if srgb is not None:
                    alp = etree.SubElement(srgb, qn('a:alpha'))
                    alp.set('val', '8000')
        except: pass

    # البطاقة الرئيسية
    s = rrect(slide, x,y,w,h, T["CD"], r_pct=8)
    # شريط علوي
    rect(slide, x,y,w,0.12, color)
    # خطوط ديكورية خفية
    for i in range(3):
        lh(slide, x+0.10, y+h-0.36-i*0.14, w*0.3, hx("FFFFFF"))

    # القيمة
    v = safe(value)
    vsz = clamp(52 - max(0,len(v)-4)*8, 26, 52)
    txt(slide, v, x+0.08, y+0.18, w-0.16, h*0.52,
        font="Calibri", size=vsz, bold=True, color=color,
        align=PP_ALIGN.CENTER)
    # خط فاصل
    lh(slide, x+0.16, y+h*0.68, w-0.32, color, 0.03)
    # التسمية
    txt(slide, safe(label), x+0.08, y+h*0.72, w-0.16, h*0.24,
        font=T["BF"], size=11, color=T["TM"], align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, safe(sub), x+0.08, y+h*0.92, w-0.16, h*0.10,
            font="Calibri", size=8, italic=True, color=T["TM"],
            align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# COVER SLIDES — 3 عائلات مختلفة جذرياً
# ═══════════════════════════════════════════════════════════════════

def cover_noir(slide, T, data):
    """غلاف NOIR — لوحة داكنة سينمائية"""
    bg(slide, T["D"])

    # ── طبقة هندسية خلفية ─────────────────────────────
    # دوائر ضخمة متداخلة
    deco_circles(slide, T, [
        (W*0.78, H*0.22, 5.8,  T["M"],  25),
        (W*0.85, H*0.55, 3.5,  T["A"],  10),
        (W*0.92, H*0.18, 2.0,  T["A"],  18),
        (-1.5,   H*0.75, 6.0,  T["M"],  22),
        (W*0.40, H*1.10, 4.5,  T["G2"], 12),
    ])
    # مثلث ديكوري
    s_tri = rect(slide, W*0.62, 0, W*0.38, H*0.08, T["A"])

    # ── شريط accent رأسي أيمن ─────────────────────────
    rect(slide, W-0.72, 0, 0.72, H, T["A"])
    rect(slide, W-0.72, 0, 0.04, H, T["G2"])

    # ── نطاق الجامعة ──────────────────────────────────
    rect(slide, 0, 0, W-0.72, 2.80, T["M"])
    lh(slide, 0, 2.80, W-0.72, T["A"], 0.06)
    lh(slide, 0, 2.86, W-0.72, T["G2"], 0.025)

    # ── رقم السنة ضخم ─────────────────────────────────
    yr = safe(data.get("year","")).replace("–","-").split("-")
    yr_txt = yr[-1].strip() if yr else ""
    if yr_txt and yr_txt.isdigit():
        txt(slide, yr_txt, 0.08, H*0.22, W*0.58, H*0.64,
            font="Calibri", size=200, bold=True,
            color=T["M"], align=PP_ALIGN.LEFT, rtl=False)

    # ── نص الجامعة والكلية ────────────────────────────
    uni = safe(data.get("university",""))
    if uni:
        txt(slide, uni, MX, 0.18, W-MX-0.90, 0.92,
            font=T["BF"], size=14, bold=True,
            color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    fac = " · ".join(filter(None,[safe(data.get("faculty","")),safe(data.get("department",""))]))
    if fac:
        txt(slide, fac, MX, 1.14, W-MX-0.90, 0.56,
            font=T["BF"], size=11, color=T["A"],
            align=PP_ALIGN.RIGHT, rtl=True)
    # مستوى الدراسة
    lvl = safe(data.get("level","ماستر 2"))
    s_pill = rrect(slide, MX, 1.82, 3.20, 0.58, T["A"], r_pct=50)
    txt(slide, "مذكرة تخرج   ·   " + lvl,
        MX, 1.82, 3.20, 0.58,
        font=T["BF"], size=11, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)

    # ── عنوان المذكرة ─────────────────────────────────
    overline(slide, safe(data.get("fieldEn","Research Thesis")),
             MX, 3.10, W-MX-0.90, T, align=PP_ALIGN.LEFT)
    lh(slide, MX, 3.44, 2.20, T["A"], 0.08)
    lh(slide, MX+2.24, 3.44, W-MX-0.90-2.24, T["M"], 0.08)
    txt(slide, safe(data.get("titleAr","")),
        MX, 3.58, W-MX-0.90, 2.90,
        font=T["BF"], size=22, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)

    # العنوان الفرنسي
    if data.get("titleFr"):
        lh(slide, MX, 6.56, W-MX-0.90, T["BO"], 0.035)
        txt(slide, safe(data.get("titleFr","")),
            MX, 6.64, W-MX-0.90, 0.56,
            font="Calibri", size=11, italic=True,
            color=T["A"], align=PP_ALIGN.LEFT, rtl=False)

    # ── شريط سفلي — طالب + مشرف ───────────────────────
    rect(slide, 0, H-0.68, W-0.72, 0.68, T["M"])
    lh(slide, 0, H-0.68, W-0.72, T["A"], 0.05)
    hw = (W-0.72-MX*2)/2
    txt(slide, "إعداد: " + safe(data.get("studentName","")),
        MX, H-0.56, hw, 0.44,
        font=T["BF"], size=12, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    txt(slide, "إشراف: " + safe(data.get("supervisor","")),
        MX+hw+0.20, H-0.56, hw-0.20, 0.44,
        font=T["BF"], size=12, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    # الكلمات المفتاحية
    kw = safe(data.get("keywords",""))
    if kw:
        rect(slide, 0, H-0.68, W-0.72, 0.68, T["M"])
        lh(slide, 0, H-0.68, W-0.72, T["A"], 0.05)
        txt(slide, "إعداد: "+safe(data.get("studentName","")),
            MX, H-0.58, hw, 0.42, font=T["BF"], size=12, bold=True,
            color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
        txt(slide, "إشراف: "+safe(data.get("supervisor","")),
            MX+hw+0.20, H-0.58, hw-0.20, 0.42, font=T["BF"], size=12, bold=True,
            color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)


def cover_vivid(slide, T, data):
    """غلاف VIVID — انقسام جريء + طاقة بصرية"""
    bg(slide, T["D"])

    # ── تقسيم رأسي جريء ───────────────────────────────
    split = W * 0.44
    rect(slide, split, 0, W-split, H, T["M"])

    # ── عناصر هندسية ──────────────────────────────────
    # مثلث ضخم يسار
    deco_circles(slide, T, [
        (split*0.2,  H*0.15, 3.8, T["M"],  30),
        (split*1.02, H*0.5,  4.2, T["A"],  18),
        (W*0.85,     H*0.85, 2.8, T["D"],  35),
        (W*0.95,     H*0.10, 1.8, T["A"],  22),
    ])

    # شريط accent رأسي فاصل
    rect(slide, split-0.06, 0, 0.12, H, T["A"])
    rect(slide, split-0.06, H*0.30, 0.12, H*0.40, T["G2"])

    # شريط accent أعلى
    rect(slide, 0, 0, W, 0.10, T["A"])

    # ── محتوى يسار (معلومات) ──────────────────────────
    # نقاط ديكورية
    decorative_grid(slide, MX, 0.22, split*0.5, 1.8, T["A"])

    txt(slide, safe(data.get("university","")),
        MX, 0.30, split-MX-0.30, 1.60,
        font=T["BF"], size=14, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    fac = " · ".join(filter(None,[safe(data.get("faculty","")),safe(data.get("department",""))]))
    if fac:
        txt(slide, fac, MX, 1.96, split-MX-0.30, 0.56,
            font=T["BF"], size=11, color=T["A"],
            align=PP_ALIGN.RIGHT, rtl=True)

    # مستوى في pill
    s = rrect(slide, MX, 2.68, min(3.0, split-MX-0.3), 0.58, T["A"], r_pct=50)
    txt(slide, safe(data.get("level","ماستر 2")),
        MX, 2.68, min(3.0, split-MX-0.3), 0.58,
        font=T["BF"], size=12, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)

    # طالب + مشرف + سنة
    txt(slide, safe(data.get("studentName","")),
        MX, H-2.10, split-MX-0.30, 0.72,
        font=T["BF"], size=18, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    lh(slide, MX, H-1.30, split-MX-0.30, T["A"], 0.04)
    txt(slide, "إشراف: " + safe(data.get("supervisor","")),
        MX, H-1.18, split-MX-0.30, 0.52,
        font=T["BF"], size=11, color=T["A"],
        align=PP_ALIGN.RIGHT, rtl=True)
    yr = safe(data.get("year",""))
    if yr:
        s2 = rrect(slide, MX, H-0.58, 2.0, 0.42, T["M"], r_pct=50)
        txt(slide, yr, MX, H-0.58, 2.0, 0.42,
            font="Calibri", size=11, bold=True,
            color=T["A"], align=PP_ALIGN.CENTER, rtl=False)

    # ── محتوى يمين (العنوان) ──────────────────────────
    rx, rw = split+0.40, W-split-0.56
    overline(slide, "عنوان المذكرة  ·  Thesis Title",
             rx, 0.24, rw, T, color=T["A"], align=PP_ALIGN.LEFT)
    lh(slide, rx, 0.58, rw, T["A"], 0.08)
    lh(slide, rx, 0.66, rw*0.60, T["G2"], 0.04)
    txt(slide, safe(data.get("titleAr","")),
        rx, 0.80, rw, 4.10,
        font=T["BF"], size=22, bold=True,
        color=T["TD"], align=PP_ALIGN.RIGHT, rtl=True)
    if data.get("titleFr"):
        lh(slide, rx, 5.02, rw, T["BO"], 0.035)
        txt(slide, safe(data.get("titleFr","")),
            rx, 5.12, rw, 0.64,
            font="Calibri", size=11, italic=True,
            color=T["D"], align=PP_ALIGN.LEFT, rtl=False)
    kw = safe(data.get("keywords",""))
    if kw:
        lh(slide, rx, H-0.72, rw, T["BO"], 0.025)
        txt(slide, "🔑 " + kw, rx, H-0.64, rw, 0.52,
            font=T["BF"], size=9, italic=True,
            color=T["D"], align=PP_ALIGN.LEFT, rtl=False)


def cover_minimal(slide, T, data):
    """غلاف MINIMAL — أبيض راقٍ + تايبوغرافي سينمائي"""
    bg(slide, hx("FAFBFE"))

    # ── شريط جانبي أيسر داكن ─────────────────────────
    rect(slide, 0, 0, 0.60, H, T["D"])
    rect(slide, 0.60, 0, 0.08, H, T["A"])

    # ── خطوط أفقية رفيعة ─────────────────────────────
    lh(slide, 0.68, H-0.10, W-0.68, T["A"], 0.10)
    lh(slide, 0.68, 0, W-0.68, T["D"], 0.08)

    # ── شبكة نقاط ديكورية ─────────────────────────────
    decorative_grid(slide, W*0.55, H*0.10, W*0.42, H*0.80, T["CE"], 4, 5)

    # ── نص رأسي في الشريط ────────────────────────────
    txt(slide, ("THESE DE MASTER  ·  " + safe(data.get("year",""))).upper(),
        0.06, 1.2, H-2.4, 0.48,
        font="Calibri", size=8, bold=True,
        color=T["TL"], align=PP_ALIGN.CENTER, rtl=False)

    # ── محتوى رئيسي ──────────────────────────────────
    cx = 0.88
    cw = W - cx - MX*0.5

    # الجامعة
    txt(slide, safe(data.get("university","")),
        cx, 0.22, cw, 1.10,
        font=T["BF"], size=14, bold=True,
        color=T["D"], align=PP_ALIGN.RIGHT, rtl=True)
    fac = " · ".join(filter(None,[safe(data.get("faculty","")),safe(data.get("department",""))]))
    if fac:
        txt(slide, fac, cx, 1.36, cw, 0.52,
            font=T["BF"], size=11, color=T["A"],
            align=PP_ALIGN.RIGHT, rtl=True)

    # خط فاصل accent بعرض متدرج
    rect(slide, cx, 2.04, 4.0, 0.10, T["A"])
    rect(slide, cx+4.0, 2.04, cw-4.0, 0.10, T["CE"])

    # عنوان ضخم
    txt(slide, safe(data.get("titleAr","")),
        cx, 2.24, cw, 3.20,
        font=T["BF"], size=26, bold=True,
        color=T["D"], align=PP_ALIGN.RIGHT, rtl=True)

    if data.get("titleFr"):
        lh(slide, cx, 5.52, cw, T["CE"], 0.04)
        txt(slide, safe(data.get("titleFr","")),
            cx, 5.62, cw, 0.64,
            font="Calibri", size=12, italic=True,
            color=T["A"], align=PP_ALIGN.LEFT, rtl=False)

    # شريط معلومات أسفل
    rect(slide, 0.68, H-1.26, W-0.68, 1.16, T["D"])
    lh(slide, 0.68, H-1.26, W-0.68, T["A"], 0.08)
    hw = (W-0.68-cx*2)/2
    txt(slide, "إعداد: " + safe(data.get("studentName","")),
        cx, H-1.10, hw, 0.52,
        font=T["BF"], size=13, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    txt(slide, "إشراف: " + safe(data.get("supervisor","")),
        cx+hw+0.20, H-1.10, hw-0.20, 0.52,
        font=T["BF"], size=13, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    # المستوى
    s = rrect(slide, W-2.60, H-0.52, 2.10, 0.36, T["A"], r_pct=50)
    txt(slide, safe(data.get("level","ماستر 2")),
        W-2.60, H-0.52, 2.10, 0.36,
        font=T["BF"], size=10, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)


def make_cover(prs, data, T):
    slide = blank(prs)
    fam = T["FAM"]
    if fam == "VIVID": cover_vivid(slide, T, data)
    elif fam == "MINIMAL": cover_minimal(slide, T, data)
    else: cover_noir(slide, T, data)
    return slide

# ═══════════════════════════════════════════════════════════════════
# SLIDE HEADER — موحد لكل الشرائح
# ═══════════════════════════════════════════════════════════════════
def slide_header(slide, T, title_ar, sub_en="", dark=True):
    """رأس الشريحة مع هندسة متقدمة"""
    fam = T["FAM"]
    h   = 1.86

    if dark:
        bg(slide, T["D"])
        # شريط accent علوي
        rect(slide, 0, 0, W, 0.10, T["A"])
        # شريط accent ثانوي
        rect(slide, 0, 0.10, W*0.35, 0.04, T["G2"])

        # ديكور ركن
        deco_circles(slide, T, [
            (W*0.82, h*0.5, 3.0, T["M"], 30),
            (W*0.96, h*1.2, 2.0, T["A"], 12),
        ])
        decorative_grid(slide, W*0.70, 0.15, W*0.28, h*0.85, T["M"])

        # overline
        overline(slide, sub_en, MX, 0.18, W-MX*2, T, align=PP_ALIGN.LEFT)
        # عنوان H1
        txt(slide, title_ar, MX, 0.44, W-MX*2, 1.06,
            font=T["HF"], size=30, bold=True,
            color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
        # خط فاصل
        rect(slide, MX, h-0.16, 2.20, 0.08, T["A"])
        rect(slide, MX+2.20, h-0.16, W-MX*2-2.20, 0.04, T["BO"])

    else:  # light
        bg(slide, hx("F8FAFF") if fam=="MINIMAL" else T["L"])
        rect(slide, 0, 0, W, 0.08, T["A"])
        rect(slide, 0, 0.08, W*0.50, 0.04, T["G2"])
        deco_circles(slide, T, [
            (-2.0, h*0.4, 4.0, T["L"], 40),
        ])
        overline(slide, sub_en, MX, 0.18, W-MX*2, T, align=PP_ALIGN.LEFT)
        txt(slide, title_ar, MX, 0.44, W-MX*2, 1.06,
            font=T["HF"], size=30, bold=True,
            color=T["TD"], align=PP_ALIGN.RIGHT, rtl=True)
        rect(slide, MX, h-0.16, 2.20, 0.08, T["A"])
        rect(slide, MX+2.20, h-0.16, W-MX*2-2.20, 0.04, T["CE"])

    return h  # content Y start

# ═══════════════════════════════════════════════════════════════════
# INTRO SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_intro(prs, data, T):
    """
    تخطيط جديد — انقسام عمودي 40/60:
    يسار: لوحة accent ملوّنة مع أيقونة ضخمة + عنوان قصير
    يمين: نص المقدمة والمقاربة مع تهوية
    """
    slide = blank(prs)
    bg(slide, T["D"])

    # ── شريط accent علوي رفيع ─────────────────────────
    rect(slide, 0, 0, W, 0.08, T["A"])
    rect(slide, 0, 0.08, W*0.30, 0.03, T["G2"])

    overview = safe(data.get("introOverview", ""))
    approach = safe(data.get("introApproach", ""))

    # ── العمود الأيسر — لوحة بصرية ───────────────────
    LW = W * 0.36
    bg_left = T["M"]
    rect(slide, 0, 0, LW, H, bg_left)
    # دائرة زخرفية ضخمة
    deco_circles(slide, T, [
        (LW * 0.50, H * 0.42, LW * 1.10, T["D"], 18),
        (LW * 0.80, H * 0.85, LW * 0.60, T["A"], 10),
    ])
    # شريط accent رأسي يفصل العمودين
    rect(slide, LW - 0.10, 0, 0.10, H, T["A"])
    rect(slide, LW - 0.14, 0, 0.04, H, T["G2"])

    # أيقونة كبيرة مركزية
    txt(slide, "📖", LW*0.10, H*0.20, LW*0.80, H*0.38,
        font="Segoe UI Emoji", size=54, align=PP_ALIGN.CENTER)
    # عنوان "مقدمة" كبير
    txt(slide, "مقدمة", LW*0.05, H*0.60, LW*0.90, H*0.24,
        font=T["HF"], size=26, bold=True,
        color=T["TL"], align=PP_ALIGN.CENTER)
    # خط زخرفي
    lh(slide, LW*0.20, H*0.86, LW*0.60, T["A"], 0.06)
    txt(slide, "INTRODUCTION", LW*0.05, H*0.89, LW*0.90, H*0.10,
        font="Calibri", size=9, bold=False, italic=True,
        color=T["A"], align=PP_ALIGN.CENTER)

    # ── العمود الأيمن — المحتوى ────────────────────────
    RX = LW + 0.32
    RW = W - RX - MX

    if overview:
        # عنوان القسم — صغير وملوّن
        s_lbl = rrect(slide, RX, 0.28, 2.60, 0.44, T["A"], r_pct=50)
        txt(slide, "نظرة عامة", RX, 0.28, 2.60, 0.44,
            font=T["BF"], size=11, bold=True,
            color=T["D"], align=PP_ALIGN.CENTER)
        # علامة اقتباس كبيرة
        txt(slide, "\u201c", RX, 0.78, 0.80, 0.80,
            font="Georgia", size=44, bold=True, color=T["A"],
            align=PP_ALIGN.LEFT, rtl=False)
        # نص المقدمة — مقتضب وبحجم مريح
        oh = H*0.46 if approach else H - 1.20
        txt(slide, overview, RX, 1.44, RW, oh,
            font=T["BF"], size=14, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)

    if approach:
        # فاصل
        lh(slide, RX, H*0.60, RW, T["BO"], 0.03)
        # label المقاربة
        s_lbl2 = rrect(slide, RX, H*0.63, 2.60, 0.44, T["M"], r_pct=50)
        txt(slide, "المقاربة النظرية", RX, H*0.63, 2.60, 0.44,
            font=T["BF"], size=11, bold=True,
            color=T["A"], align=PP_ALIGN.CENTER)
        txt(slide, approach, RX, H*0.63+0.54, RW, H*0.28,
            font=T["BF"], size=13, italic=True,
            color=T["TM"], align=PP_ALIGN.RIGHT, rtl=True)

    return slide

# ═══════════════════════════════════════════════════════════════════
# PLAN SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_plan(prs, data, T, chapters_data):
    slide = blank(prs)
    cy0   = slide_header(slide, T, "خطة الدراسة", "PLAN D'ETUDE · STUDY PLAN", dark=True)

    chs   = chapters_data[:4]
    n     = len(chs)
    if not n: return slide

    gx    = 0.22
    cw    = (W - MX*2 - gx*(n-1)) / n
    cw    = min(cw, 7.2)
    ch    = H - cy0 - 0.42

    for i, chap in enumerate(chs):
        cx  = MX + i*(cw+gx)
        sc  = T["SC"][i % len(T["SC"])]

        # بطاقة الفصل
        s = rrect(slide, cx, cy0+0.22, cw, ch, T["CD"], r_pct=8)
        if s: shadow_xml(s, blur=12, dist=4, alpha=0.18)
        # شريط علوي
        rect(slide, cx, cy0+0.22, cw, 0.09, sc)
        # شريط جانبي
        rect(slide, cx, cy0+0.22, 0.14, ch, sc)

        # رقم الفصل
        txt(slide, "F%d" % (i+1), cx+0.22, cy0+0.28, cw-0.30, 0.64,
            font="Calibri", size=24, bold=True,
            color=sc, align=PP_ALIGN.RIGHT, rtl=False)
        lh(slide, cx+0.22, cy0+0.96, cw-0.30, sc, 0.03)

        # عنوان الفصل — H3
        secs = [s for s in chap.get("sections",[]) if s][:5]
        title_h = 0.86 if secs else ch-1.10
        txt(slide, safe(chap.get("title","")),
            cx+0.22, cy0+1.02, cw-0.30, title_h,
            font=T["BF"], size=12, bold=True,
            color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)

        # المباحث
        if secs:
            sh = (H - cy0 - 1.98 - 0.42) / len(secs)
            for j, sec in enumerate(secs):
                sy = cy0 + 1.98 + j*sh
                oval(slide, cx+0.24, sy+sh*0.38, 0.10, 0.10, sc)
                txt(slide, safe(sec),
                    cx+0.40, sy+0.04, cw-0.54, sh-0.08,
                    font=T["BF"], size=10, color=T["TM"],
                    align=PP_ALIGN.RIGHT, rtl=True)
    return slide

# ═══════════════════════════════════════════════════════════════════
# PROBLEM SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_problem(prs, data, T):
    slide  = blank(prs)
    cy0    = slide_header(slide, T, "إشكالية البحث والتساؤلات",
                          "RESEARCH PROBLEM & QUESTIONS", dark=True)

    problem = safe(data.get("mainProblem",""))
    main_q  = safe(data.get("mainQuestion",""))
    subs    = [s for s in data.get("subQuestions",[]) if s][:5]

    if subs:
        lw = W*0.455 - MX - 0.10
        rw = W - MX*2 - lw - 0.28

        # ── لوحة الإشكالية (يسار) ─────────────────────
        qh = H - cy0 - (1.90 if main_q else 0.38) - 0.28
        s  = rrect(slide, MX, cy0+0.20, lw, qh, T["CD"], r_pct=8)
        if s: shadow_xml(s, blur=14, dist=5, alpha=0.20)
        rect(slide, MX, cy0+0.20, 0.10, qh, T["A"])
        txt(slide, "\u201c", MX+0.20, cy0+0.22, 1.30, 0.88,
            font="Georgia", size=46, bold=True, color=T["A"])
        txt(slide, problem, MX+0.20, cy0+0.96, lw-0.32, qh-1.08,
            font=T["BF"], size=12.5, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)

        # ── التساؤل الرئيسي ────────────────────────────
        if main_q:
            mq_y = cy0+0.20+qh+0.16
            mq_h = H - mq_y - 0.24
            s2   = rrect(slide, MX, mq_y, lw, mq_h, T["A"], r_pct=8)
            if s2: shadow_xml(s2, blur=8, dist=3, alpha=0.14)
            txt(slide, "؟", MX+0.12, mq_y+0.10, 0.76, mq_h-0.20,
                font="Georgia", size=48, bold=True,
                color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, main_q, MX+0.96, mq_y+0.14, lw-1.10, mq_h-0.28,
                font=T["BF"], size=13, bold=True,
                color=T["D"], align=PP_ALIGN.RIGHT, rtl=True)

        # ── التساؤلات الفرعية (يمين) ───────────────────
        rx    = MX + lw + 0.28
        avail = H - cy0 - 0.40
        rh    = max(1.06, (avail - 0.12*(len(subs)-1)) / len(subs))

        for i, q in enumerate(subs):
            ry = cy0+0.22 + i*(rh+0.12)
            sc = T["SC"][i % len(T["SC"])]

            s  = rrect(slide, rx, ry, rw, rh, T["CD"], r_pct=7)
            if s: shadow_xml(s, blur=8, dist=2.5, alpha=0.14)
            rect(slide, rx, ry, rw, 0.08, sc)
            # رقم دائري
            oval(slide, rx+0.18, ry+(rh-0.60)/2, 0.60, 0.60, sc)
            txt(slide, str(i+1), rx+0.18, ry+(rh-0.60)/2, 0.60, 0.60,
                font="Calibri", size=15, bold=True,
                color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, q, rx+0.90, ry+0.10, rw-1.04, rh-0.20,
                font=T["BF"], size=12, color=T["TL"],
                align=PP_ALIGN.RIGHT, rtl=True)

    else:
        qh = H - cy0 - (2.0 if main_q else 0.38) - 0.28
        s  = rrect(slide, MX, cy0+0.20, W-MX*2, qh, T["CD"], r_pct=9)
        if s: shadow_xml(s, blur=14, dist=5, alpha=0.20)
        rect(slide, MX, cy0+0.20, 0.12, qh, T["A"])
        txt(slide, "\u201c", MX+0.22, cy0+0.22, 1.60, 1.0,
            font="Georgia", size=52, bold=True, color=T["A"])
        txt(slide, problem, MX+0.22, cy0+1.08, W-MX*2-0.34, qh-1.20,
            font=T["BF"], size=14, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)
        if main_q:
            mq_y = cy0+0.20+qh+0.16
            s2   = rrect(slide, MX, mq_y, W-MX*2, H-mq_y-0.22, T["A"], r_pct=7)
            txt(slide, "؟", MX+0.14, mq_y+0.10, 0.82, H-mq_y-0.32,
                font="Georgia", size=50, bold=True, color=T["D"])
            txt(slide, main_q, MX+1.04, mq_y+0.14, W-MX*2-1.18, H-mq_y-0.32,
                font=T["BF"], size=15, bold=True, color=T["D"],
                align=PP_ALIGN.RIGHT, rtl=True)
    return slide

# ═══════════════════════════════════════════════════════════════════
# OBJECTIVES SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_objectives(prs, data, T):
    """
    تخطيط جديد — خط زمني أفقي للأهداف + لوحة جانبية للفرضيات:
    - الأهداف: أرقام دائرية بارزة فوق كل بطاقة (timeline بصري)
    - الفرضيات: لوحة داكنة يسار بتصميم مختلف كلياً عن بطاقات الأهداف
    """
    slide = blank(prs)
    dark  = (T["FAM"] != "MINIMAL")
    bg(slide, T["D"] if dark else hx("F8FAFF"))
    rect(slide, 0, 0, W, 0.08, T["A"])
    rect(slide, 0, 0.08, W*0.30, 0.03, T["G2"])

    objs  = [o for o in data.get("objectives", []) if o][:4]
    hypos = [h for h in data.get("hypotheses", []) if h][:3]

    # ── عنوان الشريحة ─────────────────────────────────
    overline(slide, "OBJECTIVES & HYPOTHESES", MX, 0.18, W-MX*2, T, align=PP_ALIGN.LEFT)
    txt(slide, "أهداف البحث والفرضيات", MX, 0.44, W-MX*2, 0.90,
        font=T["HF"], size=28, bold=True,
        color=T["TL"] if dark else T["TD"], align=PP_ALIGN.RIGHT, rtl=True)
    rect(slide, MX, 1.42, 2.0, 0.07, T["A"])
    rect(slide, MX+2.0, 1.42, W-MX*2-2.0, 0.035, T["BO"])

    cy0 = 1.60
    BODY_H = H - cy0 - 0.24

    # ── منطقة الأهداف (يمين، عرض 62%) ────────────────
    OW = W * 0.60
    n  = len(objs)

    if n:
        # خط أفقي يربط الأهداف = timeline
        line_y = cy0 + 0.52
        rect(slide, MX, line_y, OW - MX, 0.05, T["M"])

        card_w = (OW - MX - 0.14*(n-1)) / max(n, 1)
        card_h = BODY_H - 0.82

        for i, obj in enumerate(objs):
            cx = MX + i*(card_w + 0.14)
            sc = T["SC"][i % len(T["SC"])]

            # دائرة الرقم فوق الخط الزمني
            oval(slide, cx + card_w/2 - 0.44, line_y - 0.40, 0.88, 0.88, sc)
            txt(slide, str(i+1),
                cx + card_w/2 - 0.44, line_y - 0.40, 0.88, 0.88,
                font="Calibri", size=20, bold=True,
                color=T["D"] if dark else hx("FFFFFF"), align=PP_ALIGN.CENTER)

            # خط رأسي يصل الدائرة بالبطاقة
            lv(slide, cx + card_w/2 - 0.025, line_y + 0.46, 0.32, sc, 0.05)

            # البطاقة
            card_y = line_y + 0.78
            bg_c = T["CD"] if dark else hx("FFFFFF")
            s = rrect(slide, cx, card_y, card_w, card_h, bg_c, r_pct=9,
                      line_color=None if dark else T["CE"], line_w=0.8)
            if s: shadow_xml(s, blur=12, dist=4, alpha=0.16)
            rect(slide, cx, card_y, card_w, 0.09, sc)

            # أيقونة صغيرة
            icons_obj = ["🎯", "📊", "🔍", "💡"]
            txt(slide, icons_obj[i % 4],
                cx + card_w/2 - 0.30, card_y + 0.14, 0.60, 0.52,
                font="Segoe UI Emoji", size=18, align=PP_ALIGN.CENTER)

            txt(slide, safe(obj),
                cx + 0.10, card_y + 0.74, card_w - 0.20, card_h - 0.86,
                font=T["BF"], size=11.5,
                color=T["TL"] if dark else T["TD"],
                align=PP_ALIGN.CENTER, rtl=True)

    # ── منطقة الفرضيات (يسار، عرض 36%) ───────────────
    HX_start = OW + 0.20
    HW = W - HX_start - MX*0.5

    # لوحة خلفية موحدة للفرضيات
    hypo_panel = rrect(slide, HX_start, cy0 + 0.08, HW, BODY_H, T["M"], r_pct=9)
    if hypo_panel: shadow_xml(hypo_panel, blur=14, dist=5, alpha=0.20)
    rect(slide, HX_start, cy0+0.08, HW, 0.09, T["A"])

    # عنوان لوحة الفرضيات
    txt(slide, "💡", HX_start + HW/2 - 0.30, cy0 + 0.18, 0.60, 0.50,
        font="Segoe UI Emoji", size=18, align=PP_ALIGN.CENTER)
    txt(slide, "الفرضيات", HX_start + 0.10, cy0 + 0.20, HW - 0.20, 0.48,
        font=T["BF"], size=14, bold=True,
        color=T["A"], align=PP_ALIGN.CENTER)
    lh(slide, HX_start + 0.20, cy0 + 0.74, HW - 0.40, T["A"], 0.04)

    nh = len(hypos)
    if nh:
        item_h = (BODY_H - 0.90) / nh
        for i, hy in enumerate(hypos):
            iy = cy0 + 0.88 + i * item_h
            sc = T["SC"][(i+2) % len(T["SC"])]
            # رقم صغير
            s2 = rrect(slide, HX_start + 0.18, iy + item_h*0.12,
                       0.50, 0.36, sc, r_pct=50)
            txt(slide, "H%d" % (i+1),
                HX_start + 0.18, iy + item_h*0.12, 0.50, 0.36,
                font="Calibri", size=11, bold=True,
                color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, safe(hy),
                HX_start + 0.12, iy + item_h*0.12 + 0.40,
                HW - 0.24, item_h - 0.56,
                font=T["BF"], size=11,
                color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
            if i < nh - 1:
                lh(slide, HX_start+0.20, iy+item_h-0.08, HW-0.40, T["BO"], 0.02)

    return slide

# ═══════════════════════════════════════════════════════════════════
# IMPORTANCE SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_importance(prs, data, T):
    slide  = blank(prs)
    dark   = (T["FAM"] == "NOIR")
    cy0    = slide_header(slide, T, "أهمية الدراسة وأسباب اختيارها",
                          "RESEARCH SIGNIFICANCE", dark=dark)
    deco_circles(slide, T, [(-3,H*0.5,7,T["M"],20)])

    items = [x for x in data.get("importance",[]) if x]
    if data.get("reasons"): items.append(data["reasons"])
    items = items[:6]
    if not items: return slide

    n     = len(items)
    cols  = 2 if n > 3 else 1
    rows  = math.ceil(n/cols)
    gx,gy = 0.26, 0.22
    avail_w = W - MX*2
    avail_h = H - cy0 - 0.36
    cw = (avail_w - gx*(cols-1)) / cols
    ch = (avail_h - gy*(rows-1)) / rows
    icons = ["🔬","💡","📊","🎯","🌐","⚡"]

    for i, item in enumerate(items):
        col = i%cols; row = i//rows if rows>0 else 0
        col = i%cols; row = i//cols
        cx  = MX + col*(cw+gx)
        cy  = cy0+0.18 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]

        bg_c = T["CD"] if dark else T["CB"]
        s = rrect(slide, cx,cy,cw,ch, bg_c, r_pct=9,
                  line_color=None if dark else T["CE"])
        if s: shadow_xml(s, blur=10, dist=3.5, alpha=0.15)
        rect(slide, cx,cy,cw,0.09, sc)
        rect(slide, cx,cy,0.10,ch, sc)
        # أيقونة كبيرة
        txt(slide, icons[i%len(icons)],
            cx+0.20, cy+0.14, 0.80, 0.72,
            font="Segoe UI Emoji", size=24, align=PP_ALIGN.CENTER)
        # رقم دائري
        oval(slide, cx+cw-0.90, cy+0.14, 0.68, 0.68, sc)
        txt(slide, "%02d"%(i+1), cx+cw-0.90, cy+0.14, 0.68, 0.68,
            font="Calibri", size=16, bold=True,
            color=T["D"] if dark else T["CB"], align=PP_ALIGN.CENTER)
        lh(slide, cx+0.22, cy+0.96, cw-0.34, sc, 0.025)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(item), cx+0.22, cy+1.04, cw-0.34, ch-1.16,
            font=T["BF"], size=12, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide

# ═══════════════════════════════════════════════════════════════════
# METHODOLOGY SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_methodology(prs, data, T):
    slide  = blank(prs)
    dark   = (T["FAM"] != "MINIMAL")
    cy0    = slide_header(slide, T, "المنهجية والعينة والمجالات",
                          "METHODOLOGY & SAMPLE", dark=dark)

    meth    = safe(data.get("methodology",""))
    stype   = safe(data.get("sampleType",""))
    ssize   = safe(data.get("sampleSize",""))
    tool_v  = safe(data.get("tool",""))
    axes    = [a for a in data.get("toolAxes",[]) if a][:4]
    spatial = safe(data.get("spatialScope",""))
    temporal= safe(data.get("temporalScope",""))
    human_s = safe(data.get("humanScope",""))
    sw      = safe(data.get("software",""))
    tests   = [t for t in data.get("statisticalTests",[]) if t][:4]

    boxes = []
    if meth: boxes.append(("🔬","المنهج المتبع",meth))
    if stype or ssize:
        boxes.append(("👥","العينة"," · ".join(filter(None,[stype,ssize]))))
    if tool_v:
        tv = tool_v+("\n"+" · ".join(axes) if axes else "")
        boxes.append(("📋","أداة الدراسة",tv))
    if spatial or temporal or human_s:
        scope="\n".join(filter(None,[
            "📍 "+spatial if spatial else "",
            "🕐 "+temporal if temporal else "",
            "👤 "+human_s if human_s else ""]))
        boxes.append(("🌐","مجالات الدراسة",scope))
    if sw:
        swv=sw+(" · "+" · ".join(tests) if tests else "")
        boxes.append(("⚙️","البرنامج والاختبارات",swv))
    if data.get("dataSource"):
        boxes.append(("📂","مصدر البيانات",safe(data.get("dataSource",""))))

    if not boxes: return slide
    n    = len(boxes)
    cols = min(n,3)
    rows = math.ceil(n/cols)
    gx,gy = 0.24, 0.22
    bw  = (W-MX*2 - gx*(cols-1)) / cols
    bh  = (H-cy0-0.38 - gy*(rows-1)) / rows

    for i,(icon,lbl,val) in enumerate(boxes):
        col=i%cols; row=i//cols
        bx = MX+col*(bw+gx); by = cy0+0.20+row*(bh+gy)
        sc = T["SC"][i%len(T["SC"])]
        bg_c = T["CD"] if dark else T["CB"]
        s = rrect(slide, bx,by,bw,bh, bg_c, r_pct=9,
                  line_color=None if dark else T["CE"])
        if s: shadow_xml(s, blur=10, dist=3.5, alpha=0.15)
        rect(slide, bx,by,bw,0.09, sc)
        rect(slide, bx,by,0.10,bh, sc)
        txt(slide, icon, bx+0.18, by+0.12, 0.72, 0.58,
            font="Segoe UI Emoji", size=20, align=PP_ALIGN.CENTER)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, lbl, bx+0.96, by+0.14, bw-1.10, 0.50,
            font=T["BF"], size=13, bold=True, color=sc,
            align=PP_ALIGN.RIGHT, rtl=True)
        lh(slide, bx+0.18, by+0.72, bw-0.28, sc, 0.025)
        txt(slide, safe(val), bx+0.18, by+0.82, bw-0.28, bh-0.94,
            font=T["BF"], size=11.5, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide

# ═══════════════════════════════════════════════════════════════════
# KPI DASHBOARD
# ═══════════════════════════════════════════════════════════════════
def make_stats(prs, data, T):
    slide  = blank(prs)
    cy0    = slide_header(slide, T,
                          "لوحة المؤشرات الإحصائية الرئيسية",
                          "KEY PERFORMANCE INDICATORS — DASHBOARD", dark=True)
    deco_circles(slide, T, [(W*0.82,H*0.55,5,T["M"],18)])

    # تحقق وتنظيف: label وvalue كلاهما مطلوبان، والقيمة يجب ألا تتجاوز 40 حرفاً
    raw_stats = data.get("stats", [])
    stats = [
        {
            "label": str(s.get("label","")).strip()[:60],
            "value": str(s.get("value","")).strip()[:40],
            "sub":   str(s.get("sub","")).strip()[:50],
        }
        for s in raw_stats
        if str(s.get("label","")).strip() and str(s.get("value","")).strip()
    ]
    if not stats: return slide
    n    = min(len(stats),8)
    cols = min(n,4)
    rows = math.ceil(n/cols)
    gx,gy = 0.22, 0.24
    cw  = (W-MX*2 - gx*(cols-1)) / cols
    raw_ch = (H-cy0-0.38 - gy*(rows-1)) / rows
    ch  = min(raw_ch, 3.60)
    tot = rows*ch + (rows-1)*gy
    y0  = cy0+0.20 + max(0,(H-cy0-tot-0.38)/2)

    for i,s in enumerate(stats[:8]):
        col=i%cols; row=i//cols
        cx = MX+col*(cw+gx); cy = y0+row*(ch+gy)
        sc = T["SC"][i%len(T["SC"])]
        kpi_premium(slide, cx,cy,cw,ch, T, s["value"], s["label"], sc,
                    sub=s.get("sub",""))
    return slide

# ═══════════════════════════════════════════════════════════════════
# RESULTS SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_results(prs, data, T):
    """
    تخطيط جديد — رقم ضخم يسار كـ accent، النص يمين مع مساحة تنفس.
    الشرائط أفقية كاملة العرض، متناوبة بين لونين.
    """
    slide = blank(prs)
    dark  = (T["FAM"] != "MINIMAL")
    bg(slide, T["D"] if dark else hx("F8FAFF"))
    rect(slide, 0, 0, W, 0.08, T["A"])
    rect(slide, 0, 0.08, W*0.30, 0.03, T["G2"])

    results = [r for r in data.get("mainResults", []) if r][:6]
    if not results: return slide

    overline(slide, "RESEARCH FINDINGS", MX, 0.18, W-MX*2, T, align=PP_ALIGN.LEFT)
    txt(slide, "أهم نتائج البحث", MX, 0.44, W-MX*2, 0.90,
        font=T["HF"], size=28, bold=True,
        color=T["TL"] if dark else T["TD"], align=PP_ALIGN.RIGHT, rtl=True)
    rect(slide, MX, 1.42, 2.0, 0.07, T["A"])
    rect(slide, MX+2.0, 1.42, W-MX*2-2.0, 0.035, T["BO"])

    cy0   = 1.60
    n     = len(results)
    avail = H - cy0 - 0.20
    rh    = max(0.80, (avail - 0.10*(n-1)) / n)
    # عرض عمود الرقم الضخم
    NW = 1.60

    for i, res in enumerate(results):
        ry  = cy0 + i*(rh+0.10)
        sc  = T["SC"][i % len(T["SC"])]
        alt = (i % 2 == 0)
        bg_c = (T["CD"] if alt else T["M"]) if dark else (hx("FFFFFF") if alt else hx("EEF4FF"))

        # شريط كامل
        s = rrect(slide, MX, ry, W-MX*2, rh, bg_c, r_pct=6,
                  line_color=None if dark else T["CE"], line_w=0.6)
        if s: shadow_xml(s, blur=5, dist=2, alpha=0.10)

        # عمود الرقم الضخم — خلفية accent
        rect(slide, MX, ry, NW, rh, sc)
        # الرقم بخط ضخم
        num_sz = clamp(42 - max(0, n-4)*4, 28, 42)
        txt(slide, str(i+1), MX, ry, NW, rh,
            font="Calibri", size=num_sz, bold=True,
            color=T["D"], align=PP_ALIGN.CENTER)

        # فاصل رأسي خفيف
        lv(slide, MX+NW, ry+0.10, rh-0.20, T["BO"], 0.025)

        # النص مع تهوية (padding)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(res),
            MX+NW+0.22, ry+0.12, W-MX*2-NW-0.34, rh-0.24,
            font=T["BF"], size=13, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)

    return slide

# ═══════════════════════════════════════════════════════════════════
# CONCLUSION SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_conclusion(prs, data, T):
    """
    تخطيط جديد — اقتباس ضخم يملأ الشريحة مع خلفية درامية.
    شريط سفلي يحتوي مؤشرات صغيرة (إن وُجدت).
    """
    slide = blank(prs)
    bg(slide, T["D"])

    # دوائر زخرفية خلفية ضخمة
    deco_circles(slide, T, [
        (W*0.85, -2.0, 8.0,  T["M"], 20),
        (-3.0,   H*0.7, 7.0, T["M"], 15),
        (W*0.45, H*0.5, 3.5, T["A"], 8),
    ])
    decorative_grid(slide, W*0.62, H*0.05, W*0.35, H*0.55, T["M"], 4, 5)

    # شريط accent علوي
    rect(slide, 0, 0, W, 0.08, T["A"])
    rect(slide, 0, 0.08, W*0.25, 0.03, T["G2"])

    # label "الخاتمة"
    s_pill = rrect(slide, MX, 0.24, 2.80, 0.50, T["A"], r_pct=50)
    txt(slide, "✦  الخاتمة والاستنتاجات  ✦", MX, 0.24, 2.80, 0.50,
        font=T["BF"], size=11, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)

    conclusion = safe(data.get("generalConclusion", ""))

    # اقتباس ضخم مركزي
    txt(slide, "\u201c", MX, 0.92, 1.20, 1.10,
        font="Georgia", size=72, bold=True, color=T["A"],
        align=PP_ALIGN.LEFT, rtl=False)

    # لوحة شفافة خلف النص
    txt_y = 1.50
    txt_h = H - txt_y - 1.20
    s_bg = rrect(slide, MX+0.10, txt_y, W-MX*2-0.10, txt_h, T["M"], r_pct=10)
    if s_bg:
        try:
            sp = s_bg._element
            spPr = sp.find(qn('p:spPr'))
            fld = spPr.find('.//' + qn('a:solidFill'))
            if fld is not None:
                srgb = fld.find(qn('a:srgbClr'))
                if srgb is not None:
                    alp = etree.SubElement(srgb, qn('a:alpha'))
                    alp.set('val', '5000')
        except: pass
    rect(slide, MX+0.10, txt_y, 0.12, txt_h, T["A"])

    txt(slide, conclusion, MX+0.34, txt_y+0.18, W-MX*2-0.46, txt_h-0.36,
        font=T["BF"], size=15, color=T["TL"],
        align=PP_ALIGN.RIGHT, rtl=True)

    txt(slide, "\u201d", W-1.40, txt_y+txt_h-0.70, 1.20, 0.90,
        font="Georgia", size=72, bold=True, color=T["A"],
        align=PP_ALIGN.LEFT, rtl=False)

    # ── شريط سفلي أفقي ────────────────────────────────
    bar_y = H - 0.92
    rect(slide, 0, bar_y, W, 0.92, T["A"])
    rect(slide, 0, bar_y, W, 0.05, T["G2"])

    student = safe(data.get("studentName", ""))
    sup     = safe(data.get("supervisor", ""))
    yr      = safe(data.get("year", ""))

    info = "  ·  ".join(filter(None, [
        "إعداد: " + student if student else "",
        "إشراف: " + sup if sup else "",
        yr
    ]))
    txt(slide, info, MX, bar_y+0.12, W-MX*2, 0.66,
        font=T["BF"], size=12, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)
    return slide

# ═══════════════════════════════════════════════════════════════════
# RECOMMENDATIONS SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_recommendations(prs, data, T):
    """
    تخطيط جديد — خريطة طريق (roadmap):
    أرقام ضخمة متسلسلة مع أسهم اتجاه بين البطاقات.
    صفان: الصف الأول يمين، الثاني يسار (مسار متعرّج بصري).
    """
    slide = blank(prs)
    dark  = (T["FAM"] != "MINIMAL")
    bg(slide, T["D"] if dark else hx("F8FAFF"))
    rect(slide, 0, 0, W, 0.08, T["A"])
    rect(slide, 0, 0.08, W*0.30, 0.03, T["G2"])

    recs = [r for r in data.get("recommendations", []) if r][:6]
    if not recs: return slide

    overline(slide, "RECOMMENDATIONS", MX, 0.18, W-MX*2, T, align=PP_ALIGN.LEFT)
    txt(slide, "توصيات البحث", MX, 0.44, W-MX*2, 0.90,
        font=T["HF"], size=28, bold=True,
        color=T["TL"] if dark else T["TD"], align=PP_ALIGN.RIGHT, rtl=True)
    rect(slide, MX, 1.42, 2.0, 0.07, T["A"])
    rect(slide, MX+2.0, 1.42, W-MX*2-2.0, 0.035, T["BO"])

    cy0  = 1.62
    n    = len(recs)
    cols = min(n, 3)
    rows = math.ceil(n / cols)
    gx, gy = 0.20, 0.20
    cw = (W - MX*2 - gx*(cols-1)) / cols
    ch = (H - cy0 - 0.18 - gy*(rows-1)) / rows

    ICONS_REC = ["💡", "📌", "🔧", "🌐", "📈", "🤝"]

    for i, rec in enumerate(recs):
        col = i % cols
        row = i // cols
        cx  = MX + col*(cw+gx)
        cy  = cy0 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]
        bg_c = T["CD"] if dark else hx("FFFFFF")

        s = rrect(slide, cx, cy, cw, ch, bg_c, r_pct=10,
                  line_color=None if dark else T["CE"], line_w=0.8)
        if s: shadow_xml(s, blur=12, dist=4, alpha=0.18)

        # شريط علوي ملوّن
        rect(slide, cx, cy, cw, 0.10, sc)

        # الرقم الضخم — ركن يسار علوي
        num_bg = rrect(slide, cx+0.14, cy+0.18, 0.80, 0.72, sc, r_pct=8)
        txt(slide, str(i+1), cx+0.14, cy+0.18, 0.80, 0.72,
            font="Calibri", size=22, bold=True,
            color=T["D"], align=PP_ALIGN.CENTER)

        # أيقونة
        txt(slide, ICONS_REC[i % len(ICONS_REC)],
            cx+cw-0.82, cy+0.18, 0.64, 0.60,
            font="Segoe UI Emoji", size=20, align=PP_ALIGN.CENTER)

        # فاصل
        lh(slide, cx+0.14, cy+0.98, cw-0.28, sc, 0.03)

        # النص — مع تهوية جيدة وحجم مريح
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(rec),
            cx+0.14, cy+1.08, cw-0.28, ch-1.22,
            font=T["BF"], size=12, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)

        # سهم بين البطاقات في نفس الصف
        if col < cols-1 and i < n-1:
            ax = cx + cw + gx*0.15
            ay = cy + ch/2 - 0.20
            txt(slide, "←", ax, ay, gx*0.70, 0.40,
                font="Calibri", size=14, bold=True,
                color=T["A"], align=PP_ALIGN.CENTER)

    return slide

# ═══════════════════════════════════════════════════════════════════
# FUTURE SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_future(prs, data, T):
    """
    تخطيط جديد — انقسام قطري:
    النصف العلوي: خلفية داكنة مع عنوان ضخم وأيقونة صاروخ
    النصف السفلي: بطاقات أفقية مضغوطة بألوان متدرجة
    """
    slide = blank(prs)
    bg(slide, T["D"])

    items = [f for f in data.get("futureWork", []) if f][:4]
    if not items: return slide

    # ── النصف العلوي — لوحة بصرية ──────────────────────
    HEADER_H = H * 0.36

    # مثلث/قطري ديكوري
    deco_circles(slide, T, [
        (W*0.80, HEADER_H*0.5, HEADER_H*1.4, T["M"], 18),
        (W*0.95, HEADER_H*1.2, HEADER_H*0.9, T["A"],  8),
    ])
    rect(slide, 0, 0, W, 0.08, T["A"])
    rect(slide, 0, 0.08, W*0.30, 0.03, T["G2"])

    # أيقونة + عنوان
    txt(slide, "🔭", MX, HEADER_H*0.12, 1.40, HEADER_H*0.76,
        font="Segoe UI Emoji", size=46, align=PP_ALIGN.CENTER)
    overline(slide, "FUTURE RESEARCH PERSPECTIVES", MX+1.50, 0.18, W-MX*2-1.50, T, align=PP_ALIGN.LEFT)
    txt(slide, "آفاق البحث المستقبلية",
        MX+1.50, 0.46, W-MX*2-1.50, HEADER_H-0.58,
        font=T["HF"], size=28, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)

    # فاصل مائل بسيط
    lh(slide, 0, HEADER_H, W, T["A"], 0.07)
    lh(slide, 0, HEADER_H+0.07, W*0.40, T["G2"], 0.03)

    # ── البطاقات السفلية ──────────────────────────────
    cy0  = HEADER_H + 0.18
    n    = len(items)
    avail = H - cy0 - 0.18
    rh   = max(0.76, (avail - 0.12*(n-1)) / n)
    ICONS_F = ["🚀", "🌐", "💻", "🔬"]

    for i, fut in enumerate(items):
        ry  = cy0 + i*(rh+0.12)
        sc  = T["SC"][i % len(T["SC"])]
        alt = (i % 2 == 0)
        bg_c = T["CD"] if alt else T["M"]

        s = rrect(slide, MX, ry, W-MX*2, rh, bg_c, r_pct=7)
        if s: shadow_xml(s, blur=6, dist=2.5, alpha=0.12)
        rect(slide, MX, ry, W-MX*2, 0.07, sc)
        rect(slide, MX, ry, 0.10, rh, sc)

        # أيقونة
        txt(slide, ICONS_F[i % 4], MX+0.18, ry+(rh-0.54)/2, 0.54, 0.54,
            font="Segoe UI Emoji", size=18, align=PP_ALIGN.CENTER)

        # رقم badge
        s2 = rrect(slide, W-MX-0.84, ry+(rh-0.36)/2, 0.70, 0.36, sc, r_pct=50)
        txt(slide, "0%d" % (i+1), W-MX-0.84, ry+(rh-0.36)/2, 0.70, 0.36,
            font="Calibri", size=11, bold=True,
            color=T["D"], align=PP_ALIGN.CENTER)

        txt(slide, safe(fut),
            MX+0.84, ry+0.10, W-MX*2-1.76, rh-0.20,
            font=T["BF"], size=12.5, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)

    return slide

# ═══════════════════════════════════════════════════════════════════
# REFERENCES SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_references(prs, data, T):
    refs = [r for r in data.get("references",[]) if r][:6]
    if not refs: return
    slide  = blank(prs)
    dark   = (T["FAM"]=="NOIR")
    cy0    = slide_header(slide, T, "أبرز المراجع والمصادر",
                          "KEY REFERENCES", dark=dark)
    n      = len(refs)
    avail  = H - cy0 - 0.38
    rh     = max(0.90,(avail-0.12*(n-1))/n)

    for i,ref in enumerate(refs):
        ry  = cy0+0.20 + i*(rh+0.12)
        sc  = T["SC"][i%len(T["SC"])]
        alt = i%2==0
        bg_c= (T["CD"] if alt else T["M"]) if dark else (T["CB"] if alt else hx("F4F7FF"))
        s   = rrect(slide, MX,ry, W-MX*2,rh, bg_c, r_pct=6,
                    line_color=None if dark else T["CE"])
        if s: shadow_xml(s, blur=5, dist=2, alpha=0.09)
        rect(slide, MX,ry,0.10,rh, sc)
        oval(slide, MX+0.18, ry+(rh-0.58)/2, 0.58, 0.58, sc)
        txt(slide, str(i+1), MX+0.18, ry+(rh-0.58)/2, 0.58, 0.58,
            font="Calibri", size=14, bold=True,
            color=T["D"] if dark else T["CB"], align=PP_ALIGN.CENTER)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(ref), MX+0.90, ry+0.10, W-MX*2-1.02, rh-0.20,
            font=T["BF"], size=11, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide

# ═══════════════════════════════════════════════════════════════════
# THANK YOU SLIDE
# ═══════════════════════════════════════════════════════════════════
def make_final(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])

    # ── خلفية سينمائية ────────────────────────────────
    deco_circles(slide, T, [
        (W*0.65, -3.5, W*0.72, T["M"],  22),
        (-5.0,   H*0.5, W*0.8, T["A"],  10),
        (W*0.38, H*0.22, W*0.44, T["M"],18),
        (W*0.90, H*0.80, 3.0,   T["G2"],20),
    ])
    # شبكة نقاط
    decorative_grid(slide, W*0.55, H*0.05, W*0.42, H*0.60, T["M"], 4, 6)

    # شرائط accent
    rect(slide, 0, 0, W, 0.65, T["A"])
    rect(slide, 0, 0, W*0.35, 0.04, T["G2"])
    rect(slide, 0, H-0.65, W, 0.65, T["A"])
    rect(slide, W*0.65, H-0.04, W*0.35, 0.04, T["G2"])

    # ── نص الشكر — 3 لغات ─────────────────────────────
    txt(slide, "شكراً لحسن استماعكم",
        MX, H*0.20, W-MX*2, 1.70,
        font=T["HF"], size=46, bold=True,
        color=T["TL"], align=PP_ALIGN.CENTER, rtl=True)
    txt(slide, "Merci pour votre attention",
        MX, H*0.20+1.80, W-MX*2, 0.82,
        font="Calibri", size=22, italic=True,
        color=T["A"], align=PP_ALIGN.CENTER, rtl=False)
    txt(slide, "Thank you for your kind attention",
        MX, H*0.20+2.70, W-MX*2, 0.56,
        font="Calibri", size=14, italic=True,
        color=T["TM"], align=PP_ALIGN.CENTER, rtl=False)

    # ── فاصل مزخرف ────────────────────────────────────
    lh(slide, W*0.28, H*0.69, W*0.20, T["A"], 0.06)
    oval(slide, W/2-0.22, H*0.69-0.16, 0.44, 0.44, T["A"])
    lh(slide, W*0.52, H*0.69, W*0.20, T["A"], 0.06)

    # ── بيانات الطالب ─────────────────────────────────
    student  = safe(data.get("studentName",""))
    sup      = safe(data.get("supervisor",""))
    yr       = safe(data.get("year",""))
    info     = "  ·  ".join(filter(None,[
        "إعداد: "+student if student else "",
        "إشراف: "+sup if sup else ""]))
    if info:
        txt(slide, info, MX, H*0.73, W-MX*2, 0.56,
            font=T["BF"], size=13, color=T["TM"],
            align=PP_ALIGN.CENTER, rtl=True)
    univ = safe(data.get("university",""))
    if univ:
        txt(slide, univ, MX, H*0.80, W-MX*2, 0.48,
            font=T["BF"], size=11, italic=True,
            color=T["A"], align=PP_ALIGN.CENTER, rtl=True)
    if yr:
        s = rrect(slide, W/2-1.20, H-0.52, 2.40, 0.38, T["A"], r_pct=50)
        txt(slide, yr, W/2-1.20, H-0.52, 2.40, 0.38,
            font="Calibri", size=11, bold=True,
            color=T["D"], align=PP_ALIGN.CENTER, rtl=False)
    return slide

# ═══════════════════════════════════════════════════════════════════
# ORCHESTRATOR
# ═══════════════════════════════════════════════════════════════════
def generate_presentation(data: dict, output_path: str) -> None:
    key = data.get("theme","navy_gold")
    T   = PALETTES.get(key, PALETTES["navy_gold"])

    prs = Presentation()
    prs.slide_width  = Cm(W)
    prs.slide_height = Cm(H)

    cfg   = data.get("slides",{})
    def show(k): return cfg.get(k,True)
    def fl(k):   return [x for x in data.get(k,[]) if x]

    make_cover(prs,data,T)

    if show("intro") and (data.get("introOverview") or data.get("introApproach")):
        make_intro(prs,data,T)

    chs=[c for c in data.get("chapters",[]) if c.get("title")]
    if show("plan") and chs:
        make_plan(prs,data,T,chs)

    if show("problem") and (data.get("mainProblem") or data.get("mainQuestion") or fl("subQuestions")):
        make_problem(prs,data,T)

    if show("objectives") and (fl("objectives") or fl("hypotheses")):
        make_objectives(prs,data,T)

    if show("importance") and (fl("importance") or data.get("reasons")):
        make_importance(prs,data,T)

    if show("methodology") and (data.get("methodology") or data.get("sampleType") or data.get("tool")):
        make_methodology(prs,data,T)

    stats=[s for s in data.get("stats",[]) if s.get("label") and s.get("value")]
    if show("kpi") and stats:
        make_stats(prs,data,T)

    if show("results") and fl("mainResults"):
        make_results(prs,data,T)

    if show("conclusion") and data.get("generalConclusion"):
        make_conclusion(prs,data,T)

    if show("recommendations") and fl("recommendations"):
        make_recommendations(prs,data,T)

    if show("future") and fl("futureWork"):
        make_future(prs,data,T)

    if show("references") and fl("references"):
        make_references(prs,data,T)

    if show("thankyou"):
        make_final(prs,data,T)

    prs.save(output_path)
    n=len(prs.slides._sldIdLst)
    print("✅  %d slides [ultra·%s·%s] → %s"%(n,T["FAM"],key,output_path),
          file=sys.stderr)


if __name__=="__main__":
    if len(sys.argv)<3:
        print("Usage: python generator_canva.py input.json output.pptx",file=sys.stderr)
        sys.exit(1)
    with open(sys.argv[1],encoding="utf-8") as f:
        payload=json.load(f)
    generate_presentation(payload,sys.argv[2])
